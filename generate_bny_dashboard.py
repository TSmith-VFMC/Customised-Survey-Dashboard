"""
generate_bny_dashboard.py
--------------------------------
Builds a 2-page "BNY Managed Data Services - Executive Service Health"
PowerPoint deck directly from the latest VFM Cases export (.xlsx) placed in
SOURCE_DIR.

Each month the export filename changes (e.g. VFM_Cases_07-23-2026_12_30_AM.xlsx),
so this script automatically finds the most recent matching file in SOURCE_DIR
based on the timestamp encoded in the filename (falls back to file modified
time if the name can't be parsed).

Run:
    python generate_bny_dashboard.py

Output:
    BNY_Executive_Dashboard_Services_<DDMMYYYY>.pptx written to SOURCE_DIR.

NOTE ON HEURISTICS (tune in CONFIG below if your team defines these differently):
  - Repeat Incident Themes are detected automatically by extracting
    SYSTEM_TABLE_STYLE tokens (e.g. VFMC_ENTERPRISE_LOOKTHROUGH_VDM) and
    repeated phrases (e.g. "SSC SMF Fail") from the Subject column, then
    ranking by historical frequency.
Review these before distributing externally.
"""

import glob
import os
import re
from collections import Counter
from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# CONFIG
# --------------------------------------------------------------------------
CONFIG = {
    "source_dir": r"C:\Users\tsmith\Desktop\BNY Service Health",
    "source_glob": "VFM_Cases_*.xlsx",
    "closed_states": {"Closed", "Resolved", "Cancelled"},
    "priority_map": {"Critical": "P1", "High": "P2", "Moderate": "P3", "Low": "P4"},
    # Cases whose "Customer Keyword" contains this keyword are parked awaiting
    # root-cause analysis: excluded from KPIs, risk matrix, score, watchlist and
    # the priority bars, and surfaced separately (own bar + table).
    "pending_rca_column": "Customer Keyword",
    "pending_rca_keyword": "Pending RCA",
    "age_buckets": [(0, 2), (3, 5), (6, 10), (11, 10**6)],
    "age_bucket_labels": ["0-2d", "3-5d", "6-10d", "10+d"],
    # RAG per priority row x age bucket column ("G" green / "A" amber / "R" red)
    "risk_matrix": {
        "P1": ["G", "R", "R", "R"],
        "P2": ["G", "A", "R", "R"],
        "P3": ["G", "G", "A", "R"],
        "P4": ["G", "G", "G", "A"],
    },
    "aged_high_priority_days": 5,       # P2 aged threshold used for KPI + scoring
    # Overall scorecard: 100 = perfect. Each open case deducts points based on
    # its priority and how long it has been open (age bucket), per the weight
    # matrix below. P3/P4 contributions are each capped so a large low-priority
    # backlog cannot dominate the score.
    #   Score = 100 - SUM(P1) - SUM(P2) - min(SUM(P3), capP3) - min(SUM(P4), capP4)
    # Columns align with age_buckets / age_bucket_labels: 0-2d, 3-5d, 6-10d, 10+d
    "score_weights": {
        "P1": [10, 20, 25, 50],
        "P2": [1, 2, 3, 4],
        "P3": [0, 0, 2, 2],
        "P4": [0, 0, 0, 1],
    },
    "score_caps": {"P1": None, "P2": None, "P3": 20, "P4": 20},
    "attention_start": 100,
    "rag_bands": [(80, 100, "GREEN"), (50, 79, "AMBER"), (0, 49, "RED")],
    "repeat_theme_min_occurrences": 3,
    "repeat_theme_top_n": 4,
    "mttr_months": 3,          # trend shows the last 3 months...
    "mttr_month_days": 28,     # ...each a 4-week block (3 x 28d = 12 weeks total)
    "mttr_outlier_cap_days": 45,  # cap for trend chart legibility; low P1/P2
                                   # volume means a single very old backlog closure
                                   # can otherwise dominate a week's median
}

COLORS = {
    "navy": RGBColor(0x13, 0x29, 0x4B),
    "amber": RGBColor(0xF5, 0x9E, 0x0B),
    "green": RGBColor(0x16, 0xA3, 0x4A),
    "red": RGBColor(0xDC, 0x26, 0x26),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "grey_bg": RGBColor(0xF5, 0xF5, 0xF5),
    "grey_border": RGBColor(0xE6, 0xE6, 0xE6),
    "text_dark": RGBColor(0x1F, 0x29, 0x37),
    "text_muted": RGBColor(0x5B, 0x63, 0x70),
    "blue": RGBColor(0x25, 0x63, 0xEB),
    "slate": RGBColor(0x64, 0x74, 0x8B),
    "purple": RGBColor(0x7C, 0x3A, 0xED),
}

RAG_FILL = {"G": COLORS["green"], "A": COLORS["amber"], "R": COLORS["red"]}
RAG_LABEL = {"G": "Green", "A": "Amber", "R": "Red"}


# --------------------------------------------------------------------------
# DATA LOADING
# --------------------------------------------------------------------------
def find_latest_source_file(source_dir, pattern):
    candidates = glob.glob(os.path.join(source_dir, pattern))
    if not candidates:
        raise FileNotFoundError(f"No files matching {pattern!r} found in {source_dir}")

    def sort_key(path):
        name = os.path.splitext(os.path.basename(path))[0]
        m = re.search(r"VFM_Cases_(\d{2}-\d{2}-\d{4}_\d{1,2}_\d{2}_(?:AM|PM))", name)
        if m:
            try:
                return datetime.strptime(m.group(1), "%m-%d-%Y_%I_%M_%p")
            except ValueError:
                pass
        return datetime.fromtimestamp(os.path.getmtime(path))

    return max(candidates, key=sort_key)


def parse_asof_from_filename(path):
    name = os.path.splitext(os.path.basename(path))[0]
    m = re.search(r"VFM_Cases_(\d{2}-\d{2}-\d{4}_\d{1,2}_\d{2}_(?:AM|PM))", name)
    if m:
        try:
            return datetime.strptime(m.group(1), "%m-%d-%Y_%I_%M_%p")
        except ValueError:
            pass
    return datetime.fromtimestamp(os.path.getmtime(path))


def load_cases(path):
    df = pd.read_excel(path)
    df["Opened_dt"] = pd.to_datetime(df["Opened"], errors="coerce")
    df["Updated_dt"] = pd.to_datetime(df["Updated"], errors="coerce")
    df["PriorityCode"] = df["Priority"].map(CONFIG["priority_map"]).fillna("P4")
    return df


# --------------------------------------------------------------------------
# METRIC COMPUTATION
# --------------------------------------------------------------------------
def age_bucket_index(days):
    for i, (lo, hi) in enumerate(CONFIG["age_buckets"]):
        if lo <= days <= hi:
            return i
    return len(CONFIG["age_buckets"]) - 1


def detect_pending_rca(open_df):
    """Boolean Series flagging cases parked awaiting root-cause analysis.

    True where the CONFIG["pending_rca_column"] cell contains
    CONFIG["pending_rca_keyword"] (case-insensitive). The column is optional -
    older exports do not have it, in which case no case is treated as pending
    and the dashboard behaves exactly as before.
    """
    col = CONFIG["pending_rca_column"]
    if col not in open_df.columns:
        return pd.Series(False, index=open_df.index)
    values = open_df[col].fillna("").astype(str)
    return values.str.contains(CONFIG["pending_rca_keyword"], case=False, na=False, regex=False)


def detect_repeat_themes(df, open_index):
    """Auto-detect recurring incident themes from the Subject column.

    Returns a list of theme dicts (top N by historical frequency) and the set
    of currently-open case numbers that belong to any repeat theme (used both
    for the callout table and the "repeat incident" scorecard deduction).
    """
    token_pattern = re.compile(r"\b[A-Z][A-Z0-9]*(?:_[A-Z0-9]+){2,}\b")

    def extract_tokens(subject):
        return set(token_pattern.findall(str(subject)))

    all_tokens = Counter()
    token_to_rows = {}
    for idx, subj in df["Subject"].items():
        for tok in extract_tokens(subj):
            all_tokens[tok] += 1
            token_to_rows.setdefault(tok, []).append(idx)

    # also catch simple recurring phrases like "SSC SMF"
    phrase_candidates = ["SSC SMF"]
    for phrase in phrase_candidates:
        mask = df["Subject"].str.contains(phrase, case=False, na=False)
        if mask.sum() >= CONFIG["repeat_theme_min_occurrences"]:
            all_tokens[phrase] = int(mask.sum())
            token_to_rows[phrase] = list(df[mask].index)

    themes = []
    repeat_open_cases = set()
    for tok, cnt in all_tokens.most_common():
        if cnt < CONFIG["repeat_theme_min_occurrences"]:
            continue
        rows = token_to_rows[tok]
        open_rows = [r for r in rows if r in open_index]
        open_case_numbers = [df.loc[r, "Number"] for r in open_rows]
        repeat_open_cases.update(open_case_numbers)
        themes.append({
            "theme": tok.replace("_", " ").title() if "_" in tok else tok,
            "occurrences": cnt,
            "open_count": len(open_rows),
            "open_cases": open_case_numbers,
        })
        if len(themes) >= CONFIG["repeat_theme_top_n"]:
            break

    return themes, repeat_open_cases


def compute_metrics(df, asof):
    closed_states = CONFIG["closed_states"]
    open_df = df[~df["State"].isin(closed_states)].copy()
    open_df["Days Open"] = (asof - open_df["Opened_dt"]).dt.total_seconds() / 86400
    open_df["Days Open"] = open_df["Days Open"].clip(lower=0)
    open_df["AgeBucket"] = open_df["Days Open"].apply(age_bucket_index)
    open_df["RAG"] = open_df.apply(
        lambda r: CONFIG["risk_matrix"][r["PriorityCode"]][r["AgeBucket"]], axis=1
    )
    open_df["PendingRCA"] = detect_pending_rca(open_df)
    open_df = open_df.sort_values("Days Open", ascending=False)

    # Pending-RCA cases are parked awaiting root-cause analysis: they are
    # excluded from the KPIs, risk matrix, score, watchlist and priority bars,
    # and surfaced separately (own bar + table). Everything below is therefore
    # computed on the "active" (non-RCA) open cases.
    active_df = open_df[~open_df["PendingRCA"]].copy()
    rca_df = open_df[open_df["PendingRCA"]].copy()

    # ---- KPI tiles (active / non-RCA open cases) ----
    p1_open = int((active_df["PriorityCode"] == "P1").sum())
    aged_p2 = int(
        ((active_df["PriorityCode"] == "P2") & (active_df["Days Open"] > CONFIG["aged_high_priority_days"])).sum()
    )
    backlog = int(len(active_df))
    rca_pending = int(len(rca_df))

    # ---- Repeat incident themes (auto-detected, informational callout) ----
    themes, repeat_open_cases = detect_repeat_themes(df, set(open_df.index))
    repeat_incidents = int(len(repeat_open_cases))

    # ---- Priority x Age matrix counts (active cases; drive the scorecard) ----
    matrix_counts = {}
    for p in ["P1", "P2", "P3", "P4"]:
        row_counts = []
        for bucket_i in range(len(CONFIG["age_buckets"])):
            cnt = int(((active_df["PriorityCode"] == p) & (active_df["AgeBucket"] == bucket_i)).sum())
            row_counts.append(cnt)
        matrix_counts[p] = row_counts

    # ---- Management Attention Score (weighted Priority x Age model) ----
    # 100 = perfect. Each open case deducts CONFIG["score_weights"][priority][age
    # bucket] points; P3/P4 totals are capped per CONFIG["score_caps"].
    #   Score = 100 - SUM(P1) - SUM(P2) - min(SUM(P3), cap) - min(SUM(P4), cap)
    priority_names = {"P1": "P1 Critical", "P2": "P2 High", "P3": "P3 Moderate", "P4": "P4 Low"}
    score_breakdown = {}
    total_deduction = 0
    for p in ["P1", "P2", "P3", "P4"]:
        weights = CONFIG["score_weights"][p]
        counts = matrix_counts[p]
        raw = sum(wt * ct for wt, ct in zip(weights, counts))
        cap = CONFIG["score_caps"][p]
        applied = min(raw, cap) if cap is not None else raw
        score_breakdown[p] = {"raw": raw, "cap": cap, "applied": applied, "counts": counts, "weights": weights}
        total_deduction += applied

    score = max(0, min(100, CONFIG["attention_start"] - total_deduction))
    band = next(label for lo, hi, label in CONFIG["rag_bands"] if lo <= score <= hi)

    drivers = []
    for p in ["P1", "P2", "P3", "P4"]:
        b = score_breakdown[p]
        cap_note = f" (capped from {b['raw']})" if b["cap"] is not None and b["raw"] > b["cap"] else ""
        drivers.append(f"{priority_names[p]}: -{b['applied']}{cap_note}")

    # ---- Executive Watchlist (Top 5) ----
    # Priority (P1 first) > RAG severity > Days Open, so the single most
    # critical/highest-priority items always surface even if a lower
    # priority ticket happens to be older.
    rag_rank = {"R": 0, "A": 1, "G": 2}
    priority_rank = {"P1": 0, "P2": 1, "P3": 2, "P4": 3}
    watch = active_df.copy()
    watch["Eligible"] = (
        (watch["Days Open"] > 30)
        | (watch["State"] == "Escalation")
        | (watch["PriorityCode"] == "P1")
    )
    watch_pool = watch[watch["Eligible"]] if watch["Eligible"].any() else watch
    watch_pool = watch_pool.assign(
        _priority_rank=watch_pool["PriorityCode"].map(priority_rank),
        _rag_rank=watch_pool["RAG"].map(rag_rank),
    )
    watch_pool = watch_pool.sort_values(
        by=["_priority_rank", "_rag_rank", "Days Open"], ascending=[True, True, False]
    )
    top5 = watch_pool.head(5)

    # ---- Pending RCA cases (same columns as the watchlist, RCA only) ----
    rca_cases = rca_df.assign(
        _priority_rank=rca_df["PriorityCode"].map(priority_rank),
        _rag_rank=rca_df["RAG"].map(rag_rank),
    ).sort_values(by=["_priority_rank", "_rag_rank", "Days Open"], ascending=[True, True, False])

    # ---- MTTR trend (last N months) for P1 & P2 ----
    # Uses the median (not mean) resolution time per month - a handful of very
    # old tickets that were finally closed can otherwise create huge, misleading
    # spikes in a small sample. Each "month" is a 4-week block anchored on the
    # file's as-of date (3 x 28d = last 12 weeks).
    months = CONFIG["mttr_months"]
    month_days = CONFIG["mttr_month_days"]
    mttr_p1, mttr_p2, week_labels = [], [], []
    # Only genuinely resolved cases count towards MTTR - open cases (In Progress,
    # Escalation, Awaiting Info, etc.) must not be treated as "resolved".
    # NOTE: no dedicated resolved-date field exists in the export, so closure is
    # approximated by the Updated date on cases in a closed state.
    resolved = df[df["State"].isin(closed_states)].dropna(subset=["Opened_dt", "Updated_dt"]).copy()
    resolved["ResolutionDays"] = (resolved["Updated_dt"] - resolved["Opened_dt"]).dt.total_seconds() / 86400
    resolved["ResolutionDays"] = resolved["ResolutionDays"].clip(upper=CONFIG["mttr_outlier_cap_days"])
    last_p1, last_p2 = None, None
    for m_i in range(months, 0, -1):
        w_end = asof - pd.Timedelta(days=(m_i - 1) * month_days)
        w_start = w_end - pd.Timedelta(days=month_days)
        window = resolved[(resolved["Updated_dt"] > w_start) & (resolved["Updated_dt"] <= w_end)]
        p1_vals = window[window["Priority"] == "Critical"]["ResolutionDays"]
        p2_vals = window[window["Priority"] == "High"]["ResolutionDays"]
        v1 = round(p1_vals.median(), 1) if len(p1_vals) else last_p1
        v2 = round(p2_vals.median(), 1) if len(p2_vals) else last_p2
        if v1 is not None:
            last_p1 = v1
        if v2 is not None:
            last_p2 = v2
        mttr_p1.append(v1 if v1 is not None else 0)
        mttr_p2.append(v2 if v2 is not None else 0)
        # Two-line label: "Month N" plus the inclusive date range of the block,
        # so readers can see exactly which closures each point covers.
        range_start = w_end - pd.Timedelta(days=month_days - 1)
        date_range = f"{range_start.day} {range_start:%b}\u2013{w_end.day} {w_end:%b}"
        week_labels.append(f"Month {months - m_i + 1}\n{date_range}")

    open_by_priority = {
        p: int((active_df["PriorityCode"] == p).sum()) for p in ["P1", "P2", "P3", "P4"]
    }

    return {
        "asof": asof,
        "open_df": open_df,
        "p1_open": p1_open,
        "aged_p2": aged_p2,
        "repeat_incidents": repeat_incidents,
        "backlog": backlog,
        "score": score,
        "band": band,
        "drivers": drivers,
        "score_breakdown": score_breakdown,
        "total_deduction": total_deduction,
        "matrix_counts": matrix_counts,
        "top5": top5,
        "themes": themes,
        "mttr_p1": mttr_p1,
        "mttr_p2": mttr_p2,
        "week_labels": week_labels,
        "open_by_priority": open_by_priority,
        "rca_pending": rca_pending,
        "rca_cases": rca_cases,
    }


# --------------------------------------------------------------------------
# PPTX BUILDING HELPERS
# --------------------------------------------------------------------------
def add_rect(slide, left, top, width, height, fill=None, line=None):
    shp = slide.shapes.add_shape(1, left, top, width, height)  # 1 = RECTANGLE
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text, size=11, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="Calibri",
             wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font_name
        if color is not None:
            r.font.color.rgb = color
    return box


def add_kpi_tile(slide, left, top, width, height, label, value, fill, value_color=None):
    add_rect(slide, left, top, width, height, fill=fill)
    add_text(slide, left + Emu(45720), top + Emu(91440), width - Emu(91440), Inches(0.2),
              label, size=10, bold=True,
              color=COLORS["white"] if fill != COLORS["white"] else COLORS["text_dark"],
              align=PP_ALIGN.LEFT)
    add_text(slide, left + Emu(45720), top + Emu(320040), width - Emu(91440), Inches(0.3),
              str(value), size=22, bold=True,
              color=value_color or (COLORS["white"] if fill != COLORS["white"] else COLORS["text_dark"]),
              align=PP_ALIGN.LEFT)


def set_table_cell(table, r, c, text, size=10, bold=False, color=None, fill=None, align=PP_ALIGN.LEFT):
    cell = table.cell(r, c)
    cell.text = str(text)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    for p in cell.text_frame.paragraphs:
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color


def add_header(slide, title, subtitle, width):
    add_rect(slide, 0, 0, width, Inches(0.75), fill=COLORS["navy"])
    add_text(slide, Inches(0.3), Inches(0.18), width - Inches(0.6), Inches(0.4),
              title, size=24, bold=True, color=COLORS["white"])
    add_text(slide, Inches(0.3), Inches(0.55), width - Inches(0.6), Inches(0.3),
              subtitle, size=10, color=COLORS["grey_bg"])


# --------------------------------------------------------------------------
# SLIDE BUILDERS
# --------------------------------------------------------------------------
def build_slide1(prs, m):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    asof_str = m["asof"].strftime("%d %b %Y")

    add_header(slide, f"BNY Managed Data Services - Executive Service Health {asof_str}",
                f"Point-in-time daily view | As at {asof_str}", W)

    # ---- Top KPI row: Attention score tile + KPI tiles (widths fill the row) ----
    tile_top = Inches(0.95)
    tile_h = Inches(0.9)
    gap = Inches(0.1)
    left0 = Inches(0.3)

    kpis = [
        ("P1 Incidents", m["p1_open"], COLORS["red"] if m["p1_open"] > 0 else COLORS["white"]),
        ("Aged High Priority (P2)", m["aged_p2"], COLORS["amber"] if m["aged_p2"] > 0 else COLORS["white"]),
        ("Backlog (Open Cases)", m["backlog"], COLORS["white"]),
    ]
    n_tiles = len(kpis) + 1
    tile_w = int((Inches(12.73) - gap * (n_tiles - 1)) / n_tiles)

    band_fill = {"GREEN": COLORS["green"], "AMBER": COLORS["amber"], "RED": COLORS["red"]}[m["band"]]
    add_kpi_tile(slide, left0, tile_top, tile_w, tile_h, "Mgmt Attention Score",
                  f'{m["score"]}  {m["band"]}', fill=band_fill)

    for i, (label, value, fill) in enumerate(kpis):
        left = left0 + (i + 1) * (tile_w + gap)
        val_color = COLORS["white"] if fill != COLORS["white"] else COLORS["text_dark"]
        add_kpi_tile(slide, left, tile_top, tile_w, tile_h, label, value, fill=fill, value_color=val_color)

    # ---- Left panel: Priority x Age Risk Matrix ----
    panel_top = Inches(2.15)
    panel_left = Inches(0.3)
    panel_w = Inches(3.6)
    panel_h = Inches(1.9)
    add_rect(slide, panel_left, panel_top, panel_w, panel_h, fill=COLORS["white"], line=COLORS["grey_border"])
    add_text(slide, panel_left + Inches(0.1), panel_top - Inches(0.28), panel_w, Inches(0.25),
              "Priority x Age Risk Matrix", size=13, bold=True)

    rows, cols = 5, 5
    tbl_shape = slide.shapes.add_table(rows, cols, panel_left + Inches(0.1), panel_top + Inches(0.1),
                                         panel_w - Inches(0.2), panel_h - Inches(0.2))
    table = tbl_shape.table
    headers = [""] + CONFIG["age_bucket_labels"]
    for c, h in enumerate(headers):
        set_table_cell(table, 0, c, h, size=10, bold=True, fill=COLORS["grey_bg"], align=PP_ALIGN.CENTER)
    for r_i, p in enumerate(["P1", "P2", "P3", "P4"], start=1):
        set_table_cell(table, r_i, 0, p, size=10, bold=True, fill=COLORS["white"])
        for c_i in range(4):
            rag = CONFIG["risk_matrix"][p][c_i]
            count = m["matrix_counts"][p][c_i]
            text = str(count) if count else "-"
            set_table_cell(table, r_i, c_i + 1, text, size=11, bold=True,
                            color=COLORS["white"], fill=RAG_FILL[rag], align=PP_ALIGN.CENTER)

    # ---- Right panel: Executive Watchlist top 5 ----
    wl_left = Inches(4.1)
    wl_top = Inches(1.9)
    wl_w = Inches(8.9)
    add_text(slide, wl_left, wl_top, wl_w, Inches(0.25), "Executive Watchlist - Top 5", size=13, bold=True)

    wl_rows = 6
    wl_tbl_shape = slide.shapes.add_table(wl_rows, 6, wl_left, wl_top + Inches(0.3), wl_w, Inches(1.9))
    wl_table = wl_tbl_shape.table
    wl_table.columns[0].width = Inches(1.0)
    wl_table.columns[1].width = Inches(0.8)
    wl_table.columns[2].width = Inches(0.6)
    wl_table.columns[3].width = Inches(0.6)
    wl_table.columns[4].width = Inches(4.9)
    wl_table.columns[5].width = Inches(1.0)
    wl_headers = ["Case", "Priority", "Days", "RAG", "Issue Summary", "Escalated?"]
    for c, h in enumerate(wl_headers):
        set_table_cell(wl_table, 0, c, h, size=9, bold=True, fill=COLORS["grey_bg"])

    top5 = m["top5"]
    for ri in range(5):
        if ri < len(top5):
            row = top5.iloc[ri]
            rag = row["RAG"]
            set_table_cell(wl_table, ri + 1, 0, row["Number"], size=9)
            set_table_cell(wl_table, ri + 1, 1, row["Priority"], size=9)
            set_table_cell(wl_table, ri + 1, 2, f'{row["Days Open"]:.0f}', size=9, align=PP_ALIGN.CENTER)
            set_table_cell(wl_table, ri + 1, 3, RAG_LABEL[rag][0], size=9, bold=True,
                            color=COLORS["white"], fill=RAG_FILL[rag], align=PP_ALIGN.CENTER)
            subj = str(row["Subject"])
            set_table_cell(wl_table, ri + 1, 4, subj[:95] + ("..." if len(subj) > 95 else ""), size=9)
            set_table_cell(wl_table, ri + 1, 5, "Y" if row["State"] == "Escalation" else "N", size=9, align=PP_ALIGN.CENTER)
        else:
            for c in range(6):
                set_table_cell(wl_table, ri + 1, c, "", size=9)

    # ---- Repeat Incident Themes callout ----
    ri_top = Inches(4.2)
    add_text(slide, Inches(0.3), ri_top, Inches(12.7), Inches(0.25),
              "\u26A0  Repeat Incident Themes (systemic patterns requiring root-cause action, not just ticket closure)",
              size=12, bold=True, color=RGBColor(0xB4, 0x53, 0x09))

    themes = m["themes"]
    theme_rows = max(1, len(themes)) + 1
    theme_tbl_shape = slide.shapes.add_table(theme_rows, 4, Inches(0.3), ri_top + Inches(0.3), Inches(12.7), Inches(1.1))
    theme_table = theme_tbl_shape.table
    theme_table.columns[0].width = Inches(4.0)
    theme_table.columns[1].width = Inches(2.5)
    theme_table.columns[2].width = Inches(3.0)
    theme_table.columns[3].width = Inches(3.2)
    theme_headers = ["Theme (auto-detected)", "Historical Occurrences", "Currently Open", "Service Concern"]
    for c, h in enumerate(theme_headers):
        set_table_cell(theme_table, 0, c, h, size=9, bold=True, fill=COLORS["grey_bg"])
    if themes:
        for ri, theme in enumerate(themes, start=1):
            set_table_cell(theme_table, ri, 0, theme["theme"], size=9)
            set_table_cell(theme_table, ri, 1, theme["occurrences"], size=9, align=PP_ALIGN.CENTER)
            open_txt = f'{theme["open_count"]}'
            if theme["open_cases"]:
                open_txt += f' ({", ".join(theme["open_cases"][:3])})'
            set_table_cell(theme_table, ri, 2, open_txt, size=9)
            concern = ("Currently open - needs root-cause follow up" if theme["open_count"]
                       else "Recurring historically - verify preventative fix is effective")
            set_table_cell(theme_table, ri, 3, concern, size=9)
    else:
        for c in range(4):
            set_table_cell(theme_table, 1, c, "No repeat themes detected above threshold", size=9)

    # ---- Attention score driver strip ----
    driver_top = Inches(6.85)
    driver_text = f'{m["band"]} - Drivers: ' + " | ".join(m["drivers"])
    add_text(slide, Inches(0.3), driver_top, Inches(12.7), Inches(0.3), driver_text,
              size=10, bold=False, color=COLORS["text_muted"])


def build_slide2(prs, m):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    asof_str = m["asof"].strftime("%d %b %Y")

    add_header(slide, "BNY Managed Data Services - Trend, Backlog & Attention Score",
                f"As at {asof_str}", W)

    # ---- Chart 1: MTTR trend ----
    add_text(slide, Inches(0.3), Inches(0.85), Inches(6.2), Inches(0.3),
              f'Mean Time to Resolution - P1 & P2 (last {CONFIG["mttr_months"]} months, days)', size=13, bold=True)
    add_text(slide, Inches(0.3), Inches(4.65), Inches(6.2), Inches(0.25),
              f'Monthly median (4-week blocks); capped at {CONFIG["mttr_outlier_cap_days"]}d so a single legacy backlog closure does not distort the trend.',
              size=8, color=COLORS["text_muted"])
    chart_data = CategoryChartData()
    chart_data.categories = m["week_labels"]
    chart_data.add_series("P1 (days)", m["mttr_p1"])
    chart_data.add_series("P2 (days)", m["mttr_p2"])
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.3), Inches(1.2),
                                      Inches(6.2), Inches(3.4), chart_data)
    chart = gframe.chart
    chart.series[0].format.line.color.rgb = COLORS["red"]
    chart.series[1].format.line.color.rgb = COLORS["amber"]
    chart.has_legend = True

    # ---- Chart 2: Open cases by priority ----
    add_text(slide, Inches(6.8), Inches(0.85), Inches(6.2), Inches(0.3),
              f'Open Cases by Priority (as at {asof_str})', size=13, bold=True)
    bar_data = CategoryChartData()
    bar_data.categories = ["P1", "P2 High", "P3 Moderate", "P4 Low", "RCA Pending"]
    bar_data.add_series("Open Cases", [
        m["open_by_priority"]["P1"], m["open_by_priority"]["P2"],
        m["open_by_priority"]["P3"], m["open_by_priority"]["P4"], m["rca_pending"],
    ])
    bframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.8), Inches(1.2),
                                      Inches(6.2), Inches(3.4), bar_data)
    bchart = bframe.chart
    bchart.has_legend = False
    bar_fills = [COLORS["red"], COLORS["amber"], COLORS["blue"], COLORS["slate"], COLORS["purple"]]
    bpoints = bchart.plots[0].series[0].points
    for i, point in enumerate(bpoints):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = bar_fills[i]
    add_text(slide, Inches(6.8), Inches(4.65), Inches(6.2), Inches(0.25),
              "RCA Pending (purple) are open cases parked awaiting root-cause analysis - excluded from the other bars and the score.",
              size=8, color=COLORS["text_muted"])

    # ---- Management Attention Score panel ----
    # Score circle + scoring-model rulebook now live in the appendix; page 2 just
    # shows today's drivers horizontally so the RCA content below has more room.
    panel_top = Inches(4.8)
    panel_h = Inches(2.05)
    add_rect(slide, Inches(0.3), panel_top, Inches(12.7), panel_h, fill=COLORS["white"], line=COLORS["grey_border"])
    add_text(slide, Inches(0.5), panel_top + Inches(0.1), Inches(6), Inches(0.3),
              "Management Attention Score", size=16, bold=True)

    add_text(slide, Inches(0.5), panel_top + Inches(0.55), Inches(6), Inches(0.25), "Today's Drivers", size=12, bold=True)

    # Drivers laid out horizontally as pills across the panel
    dx = Inches(0.5)
    row_y = panel_top + Inches(0.9)
    for d in m["drivers"]:
        pill_w = Inches(max(1.2, 0.11 * len(d) + 0.3))
        add_rect(slide, dx, row_y, pill_w, Inches(0.4), fill=COLORS["grey_bg"], line=COLORS["grey_border"])
        add_text(slide, dx + Inches(0.1), row_y + Inches(0.04), pill_w, Inches(0.32), d, size=10)
        dx = dx + pill_w + Inches(0.2)
    total_txt = f'Score: {CONFIG["attention_start"]} - {m["total_deduction"]} = {m["score"]} ({m["band"]})'
    total_w = Inches(max(1.6, 0.11 * len(total_txt) + 0.3))
    band_fill = {"GREEN": COLORS["green"], "AMBER": COLORS["amber"], "RED": COLORS["red"]}[m["band"]]
    add_rect(slide, dx, row_y, total_w, Inches(0.4), fill=band_fill)
    add_text(slide, dx + Inches(0.1), row_y + Inches(0.04), total_w, Inches(0.32), total_txt,
              size=10, bold=True, color=COLORS["white"])

    # "Excluded from Score" note sits below the drivers row
    add_text(slide, Inches(0.5), panel_top + Inches(1.45), Inches(6), Inches(0.25), "Excluded from Score", size=11, bold=True)
    add_text(slide, Inches(0.5), panel_top + Inches(1.72), Inches(12), Inches(0.3),
              f'{m["rca_pending"]} Pending RCA case(s) are parked awaiting root-cause analysis and are listed separately on the next page.',
              size=10, color=COLORS["text_muted"])


def build_slide3(prs, m):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    asof_str = m["asof"].strftime("%d %b %Y")

    add_header(slide, "BNY Managed Data Services - Pending RCA (Root Cause Analysis in Progress)",
                f"Open cases parked awaiting RCA - excluded from the Management Attention Score, KPIs & Watchlist | As at {asof_str}", W)

    rca = m["rca_cases"]
    n = len(rca)
    add_text(slide, Inches(0.3), Inches(0.95), Inches(12.7), Inches(0.3),
              f"Pending RCA Cases ({n})", size=13, bold=True)

    rows = max(1, n) + 1
    tbl_h = Inches(min(5.8, 0.4 * rows))
    tbl_shape = slide.shapes.add_table(rows, 6, Inches(0.3), Inches(1.35), Inches(12.7), tbl_h)
    table = tbl_shape.table
    table.columns[0].width = Inches(1.3)
    table.columns[1].width = Inches(1.1)
    table.columns[2].width = Inches(0.8)
    table.columns[3].width = Inches(0.8)
    table.columns[4].width = Inches(7.3)
    table.columns[5].width = Inches(1.4)
    headers = ["Case", "Priority", "Days", "RAG", "Issue Summary", "Escalated?"]
    for c, h in enumerate(headers):
        set_table_cell(table, 0, c, h, size=10, bold=True, fill=COLORS["grey_bg"])
    if n:
        for ri in range(n):
            row = rca.iloc[ri]
            rag = row["RAG"]
            set_table_cell(table, ri + 1, 0, row["Number"], size=9)
            set_table_cell(table, ri + 1, 1, row["Priority"], size=9)
            set_table_cell(table, ri + 1, 2, f'{row["Days Open"]:.0f}', size=9, align=PP_ALIGN.CENTER)
            set_table_cell(table, ri + 1, 3, RAG_LABEL[rag][0], size=9, bold=True,
                            color=COLORS["white"], fill=RAG_FILL[rag], align=PP_ALIGN.CENTER)
            subj = str(row["Subject"])
            set_table_cell(table, ri + 1, 4, subj[:110] + ("..." if len(subj) > 110 else ""), size=9)
            set_table_cell(table, ri + 1, 5, "Y" if row["State"] == "Escalation" else "N", size=9, align=PP_ALIGN.CENTER)
    else:
        for c in range(6):
            set_table_cell(table, 1, c, "No Pending RCA cases" if c == 0 else "", size=9)


# --------------------------------------------------------------------------
# MAIN
# --------------------------------------------------------------------------
def main():
    source_path = find_latest_source_file(CONFIG["source_dir"], CONFIG["source_glob"])
    asof = pd.Timestamp(parse_asof_from_filename(source_path))
    print(f"Using source file: {source_path}")
    print(f"As-of timestamp:   {asof}")

    df = load_cases(source_path)
    m = compute_metrics(df, asof)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide1(prs, m)
    build_slide2(prs, m)
    build_slide3(prs, m)

    out_name = f'BNY_Executive_Dashboard_Services_{asof.strftime("%d%m%Y")}.pptx'
    out_path = os.path.join(CONFIG["source_dir"], out_name)
    prs.save(out_path)
    print(f"Saved: {out_path}")

    print("\n--- Summary ---")
    print(f"Mgmt Attention Score: {m['score']} ({m['band']})")
    print(f"P1 open: {m['p1_open']} | Aged P2: {m['aged_p2']} | Repeat incidents: {m['repeat_incidents']}")
    print(f"Backlog: {m['backlog']} | Pending RCA (excluded): {m['rca_pending']}")
    print("Score drivers:", "; ".join(m["drivers"]))
    print("Repeat themes:", [t["theme"] for t in m["themes"]])


if __name__ == "__main__":
    main()
